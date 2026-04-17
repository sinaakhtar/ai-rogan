import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

def extract_text_from_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        t = page.extract_text()
        if t:
            text += t + "\n"
    return text

def extract_text_from_docx(file_path: str) -> str:
    doc = Document(file_path)
    text = ""
    for para in doc.paragraphs:
        if para.text:
            text += para.text + "\n"
    return text

def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        text += run.text + " "
                text += "\n"
    return text

def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def extract_text(file_path: str) -> str:
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.txt':
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")
